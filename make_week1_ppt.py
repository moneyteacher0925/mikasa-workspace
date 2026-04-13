from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)

FONT = "에스코어 드림 7 ExtraBold"
W = RGBColor(0xFF,0xFF,0xFF)
R = RGBColor(0xFF,0x00,0x00)
G = RGBColor(0x66,0x66,0x66)
B = RGBColor(0x00,0x00,0x00)

def bg(s):
    f=s.background.fill; f.solid(); f.fore_color.rgb=B

def wm(s):
    t=s.shapes.add_textbox(Emu(10307921),Emu(6179288),Emu(1579279),Emu(369332))
    t.text_frame.word_wrap=False
    p=t.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
    r=p.add_run(); r.text="셀러 부트캠프"; r.font.name=FONT; r.font.size=Pt(14); r.font.color.rgb=G

def _run(p, text, color=W, size=Pt(60)):
    r=p.add_run(); r.text=text; r.font.name=FONT; r.font.size=size; r.font.color.rgb=color

def impact(text, highlights=None, size=Pt(60)):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    t=s.shapes.add_textbox(Emu(731520),Emu(1828800),Emu(10515600),Emu(2286000))
    t.text_frame.word_wrap=True; p=t.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    if highlights:
        for w in text.split():
            c=R if any(h in w for h in highlights) else W
            _run(p, w+" ", c, size)
    else:
        _run(p, text, W, size)
    wm(s); return s

def title_sub(title, subtitle, sub_color=W, title_size=Pt(50), sub_size=Pt(34)):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    t=s.shapes.add_textbox(Emu(731520),Emu(1500000),Emu(10515600),Emu(1200000))
    t.text_frame.word_wrap=True; p=t.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    _run(p, title, W, title_size)
    t2=s.shapes.add_textbox(Emu(731520),Emu(3200000),Emu(10515600),Emu(1500000))
    t2.text_frame.word_wrap=True; p2=t2.text_frame.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER
    _run(p2, subtitle, sub_color, sub_size)
    wm(s); return s

def title_bullets(title, bullets, title_size=Pt(44), bullet_size=Pt(28)):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    t=s.shapes.add_textbox(Emu(731520),Emu(500000),Emu(10515600),Emu(1000000))
    t.text_frame.word_wrap=True; p=t.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    _run(p, title, W, title_size)
    t2=s.shapes.add_textbox(Emu(1200000),Emu(1800000),Emu(9500000),Emu(4500000))
    t2.text_frame.word_wrap=True
    for i,b in enumerate(bullets):
        p2=t2.text_frame.paragraphs[0] if i==0 else t2.text_frame.add_paragraph()
        p2.alignment=PP_ALIGN.LEFT; p2.space_after=Pt(12)
        if isinstance(b, tuple):
            _run(p2, "• "+b[0], b[1], bullet_size)
        else:
            _run(p2, "• "+b, W, bullet_size)
    wm(s); return s

def part_slide(part_num, title):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    t=s.shapes.add_textbox(Emu(731520),Emu(1500000),Emu(10515600),Emu(1000000))
    t.text_frame.word_wrap=True; p=t.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    _run(p, f"PART {part_num}", R, Pt(40))
    t2=s.shapes.add_textbox(Emu(731520),Emu(2800000),Emu(10515600),Emu(1500000))
    t2.text_frame.word_wrap=True; p2=t2.text_frame.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER
    _run(p2, title, W, Pt(60))
    wm(s); return s

def mission(title, items):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    t=s.shapes.add_textbox(Emu(731520),Emu(500000),Emu(10515600),Emu(1000000))
    t.text_frame.word_wrap=True; p=t.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    _run(p, "🎯 MISSION", R, Pt(44))
    p2=t.text_frame.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
    _run(p2, title, W, Pt(36))
    t2=s.shapes.add_textbox(Emu(1200000),Emu(2200000),Emu(9500000),Emu(4000000))
    t2.text_frame.word_wrap=True
    for i,item in enumerate(items):
        p3=t2.text_frame.paragraphs[0] if i==0 else t2.text_frame.add_paragraph()
        p3.alignment=PP_ALIGN.LEFT; p3.space_after=Pt(14)
        _run(p3, f"☐ {item}", W, Pt(28))
    wm(s); return s

# ===== 1주차 슬라이드 생성 =====

# --- #1 오프닝: 모텔비가 아까워 ---
part_slide(1, "오프닝")
impact("모텔비가 아까워", ["모텔비가"], Pt(72))
impact("여자친구와 같이 있고 싶어도", size=Pt(48))
impact("\"일이 있다고 가봐야 했던 청년\"", ["청년"], Pt(48))
title_sub("돈이 없으면", "이렇게 됩니다", R)
impact("시간도 선택도 돈이 결정한다", ["돈이"], Pt(54))

# --- #2 오프닝: 야구선수로 실패한 인생 ---
impact("야구선수로 실패한 인생", ["실패한"], Pt(60))
title_sub("집에서는", "영양가 없는 포지션", R)
impact("뭘 해도 안 되는 느낌", size=Pt(54))
title_sub("그래서 시작했습니다", "온라인 판매", R, Pt(48), Pt(44))
impact("인생 역전의 시작", ["역전"], Pt(60))

# --- #3 동기부여 ---
part_slide(2, "동기부여")
impact("2살 많은 누나", size=Pt(60))
title_sub("명품을 손쉽게 소비", "차도 벤츠로 바꿈", R)
impact("나도 할 수 있겠다", ["할 수 있겠다"], Pt(60))
title_sub("내 목표", "벤츠 + 월 순수익 1,000만원", R, Pt(50), Pt(44))
impact("목표가 있으면 움직인다", ["목표"], Pt(54))

# --- #4 자각 ---
part_slide(3, "자각")
impact("여러분은 왜\n돈을 더 벌려고 하나요?", ["왜"], Pt(54))
title_sub("월급받으면서 계속 생활하면?", "지금 그대로 5년 뒤를 상상해보세요", W)
impact("지금 온라인 판매를 하게 된 이유는?", ["이유"], Pt(44))
title_sub("목표 설정", "온라인 판매로 얼마까지 벌 건가요?", R)
impact("목표 없이 시작하면\n반드시 포기합니다", ["반드시", "포기"], Pt(50))

# --- #5 미션: 자각미션 ---
mission("자각 미션", [
    "왜 돈을 벌어야 하나요?",
    "돈을 번다면 얼마나 벌고 싶나요?",
    "벌기 위해 어떤 걸 포기할 수 있나요?",
    "최대한 솔직하게 작성하기",
    "카톡방에 인증하기"
])
impact("솔직함이 시작입니다", ["솔직함"], Pt(54))

# --- #6 판매방식 ---
part_slide(4, "판매 방식")
impact("온라인 셀러\n판매 방식 3가지", ["3가지"], Pt(54))
title_bullets("제조 / 사입 / 위탁", [
    ("제조 = 개발자 (아인슈타인, 특허)", W),
    ("사입 = 독점권 (라이센스, 자본력)", W),
    ("위탁 = 나 잘 팔아! 리스크 NO", R),
])
impact("위탁이 가장 리스크 없고 편하다", ["위탁"], Pt(50))
title_sub("홈쇼핑도 위탁입니다", "플랫폼도 위탁 구조", W)
impact("리스크 없이 시작하세요", ["리스크 없이"], Pt(54))

# --- #7 메리트: 판매 6하원칙 (누가, 언제, 어디서) ---
part_slide(5, "판매의 본질")
impact("판매 6하원칙", ["6하원칙"], Pt(66))
title_bullets("누가 / 언제 / 어디서", [
    "내 상품을 누가 사는지 파악",
    "언제 사는지 — 시즌, 타이밍",
    "어디서 사는지 — 채널, 플랫폼",
])
impact("타겟을 모르면\n광고비만 날립니다", ["타겟"], Pt(50))
title_sub("고객을 먼저 정의하라", "그다음에 상품을 고른다", R)
impact("누가, 언제, 어디서", ["누가,", "언제,", "어디서"], Pt(60))

# --- #8 메리트: 판매 6하원칙 (무엇을, 어떻게, 왜) ---
title_bullets("무엇을 / 어떻게 / 왜", [
    "무엇을 팔아야 하는지",
    "어떻게 팔아야 하는지",
    "소비자가 왜 사야 하는지",
])
impact("왜 사야 하는지를\n설명 못하면 안 팔린다", ["왜"], Pt(50))
title_sub("6하원칙을 외워라", "판매의 모든 것이 여기에 있다", R)
impact("누가 언제 어디서\n무엇을 어떻게 왜", ["누가", "언제", "어디서", "무엇을", "어떻게", "왜"], Pt(50))

# --- #9 메리트: USP ---
impact("USP", size=Pt(80))
title_sub("Unique Selling Proposition", "내 상품의 본질은?", R, Pt(40), Pt(44))
impact("권은비의 USP는?", ["USP"], Pt(54))
title_sub("본질을 꿰뚫어야", "판매가 된다", R, Pt(50), Pt(50))
impact("한 줄로 설명 못하면\n아직 모르는 겁니다", ["한 줄"], Pt(48))

# --- #10 메리트: 이게 가장 중요해! ---
impact("이게 가장 중요해!", ["가장"], Pt(60))
impact("→ 그게 중요한 게 아니야", ["아니야"], Pt(54))
title_sub("소비자가 생각하는 중요한 포인트", "≠ 실제 구매 결정 요인", R)
impact("그걸 간파해야\n판매가 된다", ["간파"], Pt(54))
title_sub("겉으로 보이는 니즈 vs 진짜 니즈", "진짜를 찾아라", R)
impact("소비자도 자기가 왜 사는지 모른다", ["왜"], Pt(44))

# --- #11 미션: 매일 상품 1개 핵심가치 선포 ---
mission("매일 상품 1개 핵심가치 선포", [
    "세상 모든 상품 중 1개를 골라라",
    "가장 중요한 요소 1가지를 찾아라",
    "매일 카톡으로 선포하세요",
    "짧고 강하게 — 한 줄로!",
])
impact("매일 하나씩\n눈이 달라집니다", ["눈이"], Pt(54))

# --- #12 미션예시 ---
impact("미션 예시", ["예시"], Pt(60))
title_sub("에렌이 직접 보여드립니다", "상품 하나를 골라 핵심가치를 분석", W)
impact("이 상품의 핵심가치는\n단 하나입니다", ["단 하나"], Pt(48))
title_sub("핵심을 꿰뚫는 연습", "매일 하면 실력이 됩니다", R)

# --- #13 학습: 판매 수익화 구조 ---
part_slide(6, "수익화 구조")
impact("완벽하게 배워서 실행?", size=Pt(50))
impact("NO!\n구조만 알고 바로 실행", ["바로 실행"], Pt(54))
title_bullets("수익화 공식", [
    "소비자가 (판매가)",
    "- 도매가 (원가)",
    "- 수수료 (플랫폼)",
    "= 순이익 (내 돈)",
])
title_sub("마진 계산법", "판매가 - 도매가 - 수수료 = 순이익", R, Pt(44), Pt(36))
impact("계산할 줄 알면\n바로 시작할 수 있다", ["바로"], Pt(50))

# --- #14 실무: 스마트스토어 개설 ---
part_slide(7, "스마트스토어 개설")
impact("지금 바로 합니다", ["지금 바로"], Pt(66))
title_bullets("스토어 개설 순서", [
    "사업자등록",
    "스마트스토어 가입",
    "초기 세팅 (스토어명, 카테고리, 배송정보)",
    "화면 보면서 같이 진행!",
])
impact("완벽하게 만들 필요 없다", size=Pt(48))
impact("일단 열어라", ["열어라"], Pt(72))
title_sub("같이 합니다", "화면 공유하면서 실시간 진행", R)

# --- #15 실무: 위탁 도매사이트 소싱법 ---
impact("위탁 소싱법", ["소싱법"], Pt(60))
title_bullets("도매사이트 가입", [
    "도매꾹",
    "도매매",
    "오너클랜",
])
title_bullets("잘 팔리는 상품 찾는 3가지 방법", [
    ("검색량 — 사람들이 뭘 찾는지", W),
    ("트렌드 — 지금 뜨는 게 뭔지", W),
    ("시즌 — 계절/이벤트 맞춤", W),
])
title_sub("위탁 계약 프로세스", "상품 선택 → 계약 → 등록 → 판매", R)
impact("소싱이 반이다", ["반이다"], Pt(60))

# --- #16 미션: 수익화 실행 ---
mission("판매 수익화 실행 미션", [
    "수익화 구조 직접 계산해보기",
    "상품 1개 골라 마진 계산",
    "판매가 - 도매가 - 수수료 = ?",
    "카톡방에 인증!",
])
impact("계산해봐야 감이 옵니다", ["감이"], Pt(50))

# --- #17 미션: 스토어 개설 + 소싱 인증 ---
mission("스토어 개설 + 소싱 인증", [
    "스토어 개설 스크린샷 인증 (D+3)",
    "도매사이트 가입 인증",
    "상품 5개 찜하기 (D+5)",
    "카톡방에 인증!",
])
impact("실행이 답이다", ["실행"], Pt(66))
impact("1주차 끝!\n다음 주에 봅시다 🔥", ["끝!"], Pt(54))

# Save
out = "/Users/kwoneren/.openclaw/workspace/셀러부트캠프_1주차_v2.pptx"
prs.save(out)
print(f"총 슬라이드 수: {len(prs.slides)}")
print(f"저장 완료: {out}")
