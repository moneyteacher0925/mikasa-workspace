#!/usr/bin/env python3
"""셀러부트캠프 1~4주차 PPT 생성기"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import copy

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xFF, 0, 0)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_RED = RGBColor(0x99, 0, 0)

SLIDE_W = Emu(12192000)  # 16:9
SLIDE_H = Emu(6858000)

FONT_TITLE = '맑은 고딕'
FONT_BODY = '맑은 고딕'

def add_watermark(slide):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(
        Emu(SLIDE_W - Inches(2.8)), Emu(SLIDE_H - Inches(0.5)),
        Inches(2.5), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "셀러 부트캠프"
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = GRAY
    run.font.name = FONT_BODY

def set_bg_black(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg_black(slide)
    add_watermark(slide)
    return slide

def add_title_only(prs, title, font_size=44, color=WHITE, notes=""):
    slide = add_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.name = FONT_TITLE
    run.font.bold = True
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_impact(prs, text, notes=""):
    return add_title_only(prs, text, font_size=44, color=WHITE, notes=notes)

def add_impact_red(prs, text, notes=""):
    return add_title_only(prs, text, font_size=44, color=RED, notes=notes)

def add_section_divider(prs, section_name, subtitle=""):
    slide = add_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = section_name
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(40)
    run.font.color.rgb = RED
    run.font.name = FONT_TITLE
    run.font.bold = True
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(10.5), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.runs[0]
        r2.font.size = Pt(24)
        r2.font.color.rgb = WHITE
        r2.font.name = FONT_BODY
    return slide

def add_title_bullets(prs, title, bullets, notes=""):
    slide = add_slide(prs)
    # title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.color.rgb = WHITE
    run.font.name = FONT_TITLE
    run.font.bold = True
    # bullets
    txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(10.5), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        # Check if bullet starts with special markers
        if b.startswith("✅") or b.startswith("□") or b.startswith("☐"):
            p2.text = b
        else:
            p2.text = f"• {b}"
        p2.space_after = Pt(12)
        for run in p2.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = WHITE
            run.font.name = FONT_BODY
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_mission_slide(prs, title, items, notes=""):
    slide = add_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"📋 {title}"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.color.rgb = RED
    run.font.name = FONT_TITLE
    run.font.bold = True
    txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(10.5), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = f"☐ {item}"
        p2.space_after = Pt(14)
        for run in p2.runs:
            run.font.size = Pt(26)
            run.font.color.rgb = WHITE
            run.font.name = FONT_BODY
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_cover(prs, week_num, week_title):
    slide = add_slide(prs)
    # WEEK N
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11.5), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"WEEK {week_num}"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(60)
    run.font.color.rgb = RED
    run.font.name = FONT_TITLE
    run.font.bold = True
    # subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(11.5), Inches(1.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = week_title
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.runs[0]
    r2.font.size = Pt(40)
    r2.font.color.rgb = WHITE
    r2.font.name = FONT_TITLE
    r2.font.bold = True
    # 셀러 부트캠프
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(11.5), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "셀러 부트캠프"
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.runs[0]
    r3.font.size = Pt(20)
    r3.font.color.rgb = GRAY
    r3.font.name = FONT_BODY

def create_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

# ============================================================
# WEEK 1
# ============================================================
def gen_week1():
    prs = create_prs()
    add_cover(prs, 1, "위탁판매 시작하기")

    # --- 1. 오프닝 ---
    add_section_divider(prs, "OPENING", "오프닝")

    # #1 모텔비가 아까워
    add_impact(prs, "모텔비가 아까워", "여자친구와 같이 있고싶어도 일이 있다고 가봐야 했던 시절 이야기로 시작합니다.")
    add_title_bullets(prs, "돈이 없으면 이렇게 됩니다", [
        "여자친구와 같이 있고 싶어도...",
        "일이 있다고 거짓말하고 가야 했던 청년",
        "모텔비조차 아까웠던 그 시절",
    ], "돈이 없어서 겪었던 실제 에피소드를 솔직하게 공유합니다.")
    add_impact_red(prs, "돈이 없으면\n선택지가 없다", "핵심 메시지를 강하게 전달합니다.")

    # #2 야구선수로 실패한 인생
    add_impact(prs, "야구선수로 실패한 인생", "야구를 그만두고 나서의 방황기를 이야기합니다.")
    add_title_bullets(prs, "온라인 판매를 시작하게 된 이유", [
        "집에서는 영양가 없는 포지션",
        "할 줄 아는 게 없었던 현실",
        "그래서 온라인 판매를 시작했다",
    ], "야구를 그만둔 후 온라인 판매를 시작하게 된 배경을 공유합니다.")

    # --- 3. 동기부여 ---
    add_section_divider(prs, "동기부여", "목표를 세우는 이유")

    add_impact(prs, "2살 많은 누나가\n명품을 손쉽게 소비하고\n차도 벤츠로 바꾸더라", "주변에서 받은 동기부여 사례를 공유합니다.")
    add_title_bullets(prs, "내 목표는 벤츠, 월 순수익 1,000만원", [
        "명품을 손쉽게 소비하는 사람들",
        "나도 저렇게 될 수 있을까?",
        "목표: 월 순수익 1,000만원",
        "벤츠 = 내 동기부여의 상징",
    ], "구체적인 목표가 있어야 움직인다는 메시지를 전달합니다.")

    # --- 4. 자각 ---
    add_section_divider(prs, "자각", "왜 돈을 벌어야 하는가?")

    add_impact_red(prs, "여러분은 왜\n돈을 더 벌려고 하나요?", "수강생에게 직접 질문을 던집니다.")
    add_title_bullets(prs, "목표설정", [
        "월급받으면서 계속 생활하면?",
        "지금 온라인 판매를 하게 된 이유는?",
        "온라인 판매로 얼마까지 벌 건가요?",
        "구체적 숫자로 목표를 세워라",
    ], "목표 없이 시작하면 중간에 무조건 포기합니다. 지금 바로 적어보세요.")

    # --- 5. 미션1 ---
    add_section_divider(prs, "MISSION 1", "자각 미션")

    add_impact(prs, "자각미션", "첫 번째 미션을 안내합니다.")
    add_mission_slide(prs, "1주차 미션 1번", [
        "왜 돈을 벌어야 하나요?",
        "돈을 번다면 얼마나 벌고 싶나요?",
        "벌기 위해 어떤 걸 포기할 건가요?",
        "최대한 솔직하게 작성하세요",
    ], "진심을 담아 작성해야 합니다. 대충 쓰면 대충 벌게 됩니다.")

    # --- 6. 판매방식 ---
    add_section_divider(prs, "판매방식", "온라인 셀러의 3가지 길")

    add_impact(prs, "온라인 셀러\n판매방식 3가지", "세 가지 판매 방식을 비교합니다.")
    add_title_bullets(prs, "제조 vs 사입 vs 위탁", [
        "제조: 개발자 마인드 (특허, 아인슈타인)",
        "사입: 독점권 (라이센스, 자본력)",
        "위탁: 나 잘 팔아! 리스크 지기 싫어",
    ], "세 가지 방식의 차이를 명확히 설명합니다.")
    add_impact_red(prs, "위탁이\n가장 리스크 없고\n편하다", "위탁판매의 장점을 강조합니다.")
    add_title_bullets(prs, "위탁판매의 본질", [
        "홈쇼핑, 플랫폼 = 위탁의 대표 사례",
        "재고 부담 ZERO",
        "판매와 마케팅에만 집중",
        "리스크 없이 시작 가능",
    ], "위탁판매가 초보자에게 가장 적합한 이유를 설명합니다.")

    # --- 7. 판매 6하원칙 (누가, 언제, 어디서) ---
    add_section_divider(prs, "판매 6하원칙", "누가, 언제, 어디서")

    add_impact(prs, "누가 사는지\n언제 사는지\n어디서 사는지", "6하원칙 전반부를 설명합니다.")
    add_title_bullets(prs, "판매 6하원칙 (1/2)", [
        "누가: 내 상품의 타겟 고객은?",
        "언제: 시즌? 시간대? 상황?",
        "어디서: 어떤 플랫폼? 어떤 채널?",
        "이 3가지를 파악해야 판매가 된다",
    ], "타겟, 타이밍, 채널을 먼저 파악해야 합니다.")

    # --- 8. 판매 6하원칙 (무엇을, 어떻게, 왜) ---
    add_impact(prs, "무엇을 팔아야 하는지\n어떻게 팔아야 하는지\n왜 사야 하는지", "6하원칙 후반부입니다.")
    add_title_bullets(prs, "판매 6하원칙 (2/2)", [
        "무엇을: 어떤 상품을 선택할 것인가",
        "어떻게: 마케팅 방법, 가격 전략",
        "왜: 소비자가 내 상품을 사야 하는 이유",
        "6가지를 모두 답할 수 있어야 한다",
    ], "6하원칙 후반부 — 상품, 방법, 이유를 설명합니다.")

    # --- 9. USP ---
    add_section_divider(prs, "USP", "나만의 차별점")

    add_impact_red(prs, "내 상품의 본질은\n무엇인가?", "USP 개념을 소개합니다.")
    add_title_bullets(prs, "나만의 차별점 만들기 (USP)", [
        "USP = Unique Selling Proposition",
        "내 상품만의 유일한 가치",
        "권은비의 USP는?",
        "차별점이 없으면 가격으로만 싸운다",
    ], "USP가 없으면 레드오션에서 가격 경쟁만 하게 됩니다.")

    # --- 10. 소비자 간파 ---
    add_impact(prs, "이게 가장 중요해!\n→ 그게 중요한 게 아니야", "소비자가 생각하는 것과 실제 구매 결정 요인의 차이를 설명합니다.")
    add_title_bullets(prs, "소비자 간파하기", [
        "소비자가 '중요하다'고 말하는 것 ≠ 실제 구매 결정 요인",
        "표면적 니즈 vs 진짜 니즈",
        "그걸 간파해야 판매가 된다",
        "예시를 통해 직접 체험",
    ], "소비자의 진짜 구매 동기를 파악하는 훈련입니다.")

    # --- 11. 미션2 ---
    add_section_divider(prs, "MISSION 2", "매일 핵심가치 선포")

    add_impact(prs, "매일 상품 1개\n핵심가치 선포", "두 번째 미션을 안내합니다.")
    add_mission_slide(prs, "1주차 미션 2번", [
        "세상의 모든 상품 중 1개를 골라라",
        "그 상품의 가장 중요한 요소 1가지를 찾아라",
        "매일 카톡으로 선포하세요",
        "이 훈련이 판매력의 기본이 됩니다",
    ], "매일 꾸준히 해야 하는 미션입니다. 하루도 빠지지 마세요.")

    # --- 12. 미션예시 ---
    add_title_bullets(prs, "돈되는 카피라이팅 — 예시", [
        "에렌이 직접 상품 하나를 골라 시연",
        "핵심가치를 어떻게 뽑아내는지",
        "한 줄로 표현하는 연습",
    ], "강사가 직접 시연하면서 방법을 보여줍니다.")

    # --- 13. 학습: 수익화 구조 ---
    add_section_divider(prs, "학습", "판매 수익화 구조")

    add_impact_red(prs, "완벽하게 배워서 실행?\nNO!\n구조만 알고 바로 실행!", "완벽주의를 버리고 바로 실행하라는 메시지입니다.")
    add_title_bullets(prs, "돈되는 판매 수익화 구조", [
        "소비자가 → 도매가 → 수수료 → 순이익",
        "마진 계산법을 숫자로 이해",
        "구조를 알면 바로 실행 가능",
        "완벽하게 배우려 하면 시작도 못 한다",
    ], "수익 구조를 간단히 설명하고 바로 실행으로 넘어갑니다.")

    # --- 14. 실무: 스마트스토어 ---
    add_section_divider(prs, "실무", "스마트스토어 개설 실습")

    add_impact(prs, "스마트스토어 개설을\n지금 바로 합니다", "실습을 시작합니다.")
    add_title_bullets(prs, "스마트스토어 개설 실습", [
        "사업자등록 → 스토어 가입",
        "초기 세팅: 스토어명, 카테고리, 배송정보",
        "화면 보면서 같이 진행",
        "완벽하게 만들 필요 없다. 일단 열어라!",
    ], "함께 화면을 보면서 스마트스토어를 개설합니다.")
    add_impact_red(prs, "완벽하게 만들 필요 없다\n일단 열어라!", "실행이 중요하다는 메시지를 한 번 더 강조합니다.")

    # --- 15. 실무: 소싱 ---
    add_impact(prs, "위탁 도매사이트 소싱법", "소싱 방법을 알려줍니다.")
    add_title_bullets(prs, "도매사이트 소싱법", [
        "도매꾹 / 도매매 / 오너클랜 가입",
        "잘 팔리는 상품 찾는 3가지 방법:",
        "  1) 검색량 기반",
        "  2) 트렌드 기반",
        "  3) 시즌 기반",
    ], "도매 사이트 가입부터 소싱까지 같이 진행합니다.")
    add_title_bullets(prs, "위탁 계약 프로세스", [
        "도매처에 위탁 요청",
        "판매가 / 도매가 / 배송비 확인",
        "주문 발생 시 도매처가 직배송",
        "재고 리스크 ZERO",
    ], "위탁 계약이 어떻게 진행되는지 설명합니다.")

    # --- 16. 미션3 ---
    add_section_divider(prs, "MISSION 3", "수익화 실행 미션")

    add_mission_slide(prs, "1주차 미션 3번", [
        "수익화 구조를 바로 실행하기",
        "배운 마진 계산법으로 직접 계산",
        "상품 1개 선정 → 마진 시뮬레이션",
    ], "바로 실행하는 것이 핵심입니다.")

    # --- 17. 미션 추가 ---
    add_mission_slide(prs, "1주차 추가 미션", [
        "스토어 개설 스크린샷 인증 (D+3)",
        "도매사이트 가입 + 상품 5개 찜 (D+5)",
    ], "기한 내에 반드시 완료해주세요.")

    # --- 미션 요약 ---
    add_section_divider(prs, "1주차 미션 요약", "")
    add_mission_slide(prs, "1주차 전체 미션 정리", [
        "미션1: 자각미션 — 왜 돈을 벌어야 하는지 솔직하게 작성",
        "미션2: 매일 상품 1개 핵심가치 선포 (카톡)",
        "미션3: 수익화 구조 실행 — 마진 시뮬레이션",
        "추가: 스토어 개설 인증 (D+3)",
        "추가: 도매사이트 가입 + 상품 5개 찜 (D+5)",
    ], "모든 미션을 정리하여 다시 한번 안내합니다.")

    add_impact(prs, "다음 주에 만나요! 🔥", "1주차를 마무리합니다.")

    prs.save('/Users/kwoneren/.openclaw/workspace/셀러부트캠프_1주차.pptx')
    print(f"Week 1: {len(prs.slides)} slides")

# ============================================================
# WEEK 2
# ============================================================
def gen_week2():
    prs = create_prs()
    add_cover(prs, 2, "사고력 트레이닝 & 상품 등록")

    # 1. 오프닝/복습
    add_section_divider(prs, "OPENING", "복습 & 동기부여")

    add_impact(prs, "얼마를 벌어야\n자유롭게 살 수 있을까?", "돈 버는 이유를 다시 짚어봅니다.")
    add_title_bullets(prs, "돈 버는 이유 다시 짚기", [
        "자유로운 삶을 위한 최소 금액은?",
        "월급만으로는 충분한가?",
        "지금 이 자리에 앉아있는 이유를 떠올려라",
    ], "1주차에 세운 목표를 다시 상기시킵니다.")

    # 2. 모텔→호텔
    add_impact_red(prs, "신라모텔에서\n신라호텔로", "생활 수준의 변화를 시각적으로 보여줍니다.")
    add_title_bullets(prs, "월 1,000만원 벌면 벌어지는 현상", [
        "숙소가 모텔에서 호텔로",
        "선택지가 늘어난다",
        "돈이 자유를 만든다",
    ], "구체적인 변화를 이야기하여 동기를 부여합니다.")

    # 3. 미션피드백
    add_section_divider(prs, "미션 피드백", "1주차 미션 1번")
    add_title_bullets(prs, "1주차 미션1번 피드백", [
        "수강생별 목표 리뷰",
        "잘 쓴 예시 공유",
        "진심 담긴 글 vs 대충 쓴 글 차이",
        "솔직함이 성과를 만든다",
    ], "수강생들의 자각미션을 피드백합니다.")

    # 4. 판매방식
    add_impact(prs, "위탁판매 =\n판매 + 마케팅\n잘해야 한다", "위탁의 핵심 역량을 강조합니다.")
    add_title_bullets(prs, "리스크 없이 월 1억까지", [
        "위탁판매의 스케일링 가능성",
        "재고 없이 월 1억 매출 구조",
        "핵심은 판매력과 마케팅력",
    ], "위탁판매로도 큰 매출이 가능한 구조를 설명합니다.")

    # 5-6. 복습
    add_section_divider(prs, "복습", "판매 6하원칙")
    add_title_bullets(prs, "(복습) 판매 6하원칙 — 누가, 언제, 어디서", [
        "지난주 배운 6하원칙을",
        "수강생 실제 상품에 적용",
        "누가: 타겟 재정의",
        "언제, 어디서: 채널과 타이밍",
    ], "이론이 아닌 실제 상품에 적용하는 복습입니다.")
    add_title_bullets(prs, "(복습) 판매 6하원칙 — 무엇을, 어떻게, 왜", [
        "소비자의 욕구를 파악하고",
        "내 상품과 연결하는 훈련",
        "왜 사야 하는지를 한 문장으로",
    ], "6하원칙 후반부를 실전 적용합니다.")

    # 7. 미션피드백2
    add_section_divider(prs, "미션 피드백", "1주차 미션 2번")
    add_title_bullets(prs, "1주차 미션2번 피드백", [
        "수강생 제출물 분석",
        "핵심 꿰뚫은 사례 vs 겉핥기",
        "좋은 분석의 공통점",
    ], "핵심가치 선포 미션의 피드백입니다.")

    # 8. 총량의 저주
    add_section_divider(prs, "자각", "집중의 힘")
    add_impact_red(prs, "총량의 저주", "너무 많은 일을 하면 집중할 수 없다는 메시지입니다.")
    add_title_bullets(prs, "판매 전문가가 되기", [
        "너무 많은 일 = 집중 불가",
        "김치찌개 전문점 vs 분식집",
        "하나에 집중해야 전문가가 된다",
        "판매 전문가가 되어라",
    ], "선택과 집중의 중요성을 강조합니다.")

    # 9. 동기부여
    add_section_divider(prs, "동기부여", "")
    add_impact(prs, "동기부여(기름) 없으면\n앞으로 못 간다", "동기부여의 중요성을 이야기합니다.")
    add_title_bullets(prs, "월천 이상 번 수강생들의 공통점", [
        "목적이 명확하다",
        "동기부여가 꾸준하다",
        "기름이 떨어지면 차도 멈춘다",
        "여러분의 기름은 무엇인가?",
    ], "성공한 수강생들의 공통 패턴을 공유합니다.")

    # 10. 지피지기
    add_section_divider(prs, "사고력 트레이닝", "")
    add_impact(prs, "지피지기 백전불태", "손자병법의 핵심을 판매에 적용합니다.")
    add_title_bullets(prs, "지피지기 백전불태", [
        "적의 사정을 알고 나의 사정을 알면",
        "백번 싸워도 위태롭지 않다",
        "경쟁자를 알고, 소비자를 알고, 나를 알아라",
    ], "시장 분석의 중요성을 손자병법으로 풀어냅니다.")

    # 11. 통찰 4단계
    add_impact_red(prs, "통찰 4단계\n실전 적용", "사고력 트레이닝의 핵심입니다.")
    add_title_bullets(prs, "통찰 4단계 실전 적용", [
        "1인칭 → 전지적 참견 시점으로 전환",
        "수강생 찜한 상품 5개로 즉석 실습",
        "이 상품 사는 사람은 누구인가?",
        "뭘 검색하고 왜 사는가?",
    ], "수강생의 실제 상품으로 즉석 실습합니다.")

    # 12. 다이어트약 예시
    add_title_bullets(prs, "맞춤형 상품 제안 — 다이어트약", [
        "타겟A: 뚱뚱 + 게으른 → 편한 방법 강조",
        "타겟B: 결혼식 전 단기간 → 빠른 효과 강조",
        "같은 상품도 타겟에 따라 마케팅이 달라진다",
    ], "타겟에 따라 같은 상품의 소구점이 달라지는 예시입니다.")

    # 13. 마인드셋
    add_section_divider(prs, "인지", "마인드셋의 차이")
    add_impact(prs, "똑같은 방법을 알려줘도\n수익이 다른 이유", "마인드셋의 중요성을 설명합니다.")
    add_title_bullets(prs, "마인드셋이 실무보다 중요한 이유", [
        "같은 방법 → 월 300 vs 월 3,000 vs 월 3억",
        "차이는 방법이 아니라 사람",
        "여러분 자체를 바꿔야 성과가 바뀐다",
    ], "실무 스킬보다 마인드셋이 먼저라는 메시지를 전달합니다.")

    # 14. 키워드 리서치
    add_section_divider(prs, "실무", "키워드 리서치")
    add_impact(prs, "돈되는 키워드 찾기", "키워드 리서치 실습을 시작합니다.")
    add_title_bullets(prs, "키워드 리서치 실습", [
        "네이버 검색광고 키워드도구 사용법",
        "검색량 · 경쟁강도 보는 법",
        "블루오션 키워드 기준 3가지",
        "6하원칙의 '누가'를 키워드로 바꾸는 훈련",
    ], "실제 도구를 사용하면서 키워드를 찾는 방법을 알려줍니다.")

    # 15. 상품 등록
    add_impact(prs, "스마트스토어\n상품 등록 A to Z", "상품 등록 실습입니다.")
    add_title_bullets(prs, "상품 등록 실습", [
        "카테고리 → 상품명 (상위노출 공식)",
        "옵션 → 가격 (마진 역산)",
        "대표이미지 설정",
        "화면 보면서 같이 등록",
    ], "오늘 1개는 등록하고 간다는 목표로 진행합니다.")
    add_impact_red(prs, "오늘 1개는\n등록하고 간다", "실습 목표를 명확히 합니다.")

    # 16. 경쟁 상품 분석
    add_title_bullets(prs, "경쟁 상품 분석", [
        "상위노출 상품 5개 뜯어보기",
        "리뷰수 / 가격 / 배송 분석",
        "상세페이지 / 키워드 분석",
        "내가 이길 수 있는 포인트 찾기",
    ], "경쟁자를 분석해서 내 차별점을 만듭니다.")

    # 17. 미션
    add_section_divider(prs, "2주차 미션", "")
    add_mission_slide(prs, "2주차 미션 안내", [
        "상품 1개 등록 완료 스크린샷 (D+3)",
        "경쟁상품 3개 분석표 작성 (D+5)",
        "매일 핵심가치 선포 계속",
    ], "2주차 미션을 안내합니다.")

    add_impact(prs, "다음 주에 만나요! 🔥", "2주차를 마무리합니다.")

    prs.save('/Users/kwoneren/.openclaw/workspace/셀러부트캠프_2주차.pptx')
    print(f"Week 2: {len(prs.slides)} slides")

# ============================================================
# WEEK 3
# ============================================================
def gen_week3():
    prs = create_prs()
    add_cover(prs, 3, "마케팅 실전 & 상세페이지")

    # 1. 동기부여
    add_section_divider(prs, "동기부여", "어라? 이게 되네?")
    add_impact(prs, "되고 싶은 대상들의 삶", "동기부여로 시작합니다.")
    add_title_bullets(prs, "어라? 이게 되네?", [
        "사람들은 될 것 같을 때 더 한다",
        "다이어트 조금 빠지면 → 더 열심히",
        "얼굴 이쁘네 들으면 → 더 꾸민다",
        "실행력 UP의 비밀 = 작은 성공 경험",
    ], "작은 성공이 실행력을 만든다는 메시지입니다.")

    # 2. 트레이닝 - 썸네일
    add_section_divider(prs, "트레이닝", "첫인상 꾸미기")
    add_impact_red(prs, "첫인상이 전부다\n소개팅 5초 판단", "첫인상의 중요성을 강조합니다.")
    add_title_bullets(prs, "내 상품 첫인상 꾸미기", [
        "한눈에 들어오는 썸네일",
        "상세페이지, 콘텐츠",
        "월드컵 이상형 방식으로 실습",
        "소비자는 5초 안에 판단한다",
    ], "썸네일과 상페의 첫인상을 개선하는 훈련입니다.")

    # 3. 좋은 판매자
    add_impact(prs, "똑똑한 놈이\n세상을 지배한다", "전략적 사고의 중요성입니다.")
    add_title_bullets(prs, "좋은 판매자가 되는 법", [
        "감독님의 새벽 2시 트레이닝 이유",
        "학부모 마케팅: 부모님이 돈 주면 4번 타자",
        "당신의 잘못이 아닙니다",
        "전략적으로 생각하라",
    ], "판매에서도 전략적 사고가 필수라는 메시지입니다.")

    # 4. 미션피드백
    add_section_divider(prs, "미션 피드백", "2주차 미션")
    add_title_bullets(prs, "2주차 미션 피드백", [
        "수강생 등록 상품 화면 공유 + 즉석 피드백",
        "상품명 / 가격 / 카테고리 수정 포인트",
        "잘한 사례 vs 아쉬운 사례 비교",
    ], "수강생의 상품 등록 결과를 즉석 피드백합니다.")

    # 5. 소비자 구매과정
    add_section_divider(prs, "학습", "소비자 구매 여정")
    add_impact(prs, "소비자들은 어떻게\n내 상품까지 오게 될까?", "소비자 구매과정을 설명합니다.")
    add_title_bullets(prs, "소비자 구매과정 6단계", [
        "1. 현상 — 문제 인식",
        "2. 욕구 — 해결하고 싶다",
        "3. 계획 — 정보 탐색",
        "4. 행동 — 검색/비교",
        "5. 대안 — 후보군 비교",
        "6. 결과 — 구매 결정",
    ], "소비자가 구매에 이르는 6단계를 설명합니다.")

    # 6. 관여도
    add_impact_red(prs, "다 똑같은 상품이\n아닙니다", "관여도 개념을 소개합니다.")
    add_title_bullets(prs, "상품 2가지의 분류 — 관여도", [
        "고관여 상품: 비싸고 신중하게 구매 (가전, 자동차)",
        "저관여 상품: 싸고 즉흥적 구매 (간식, 생필품)",
        "관여도에 따라 소비 과정이 다르다",
        "구매 전 대부분 결정되어 있음",
    ], "관여도에 따른 소비자 행동 차이를 설명합니다.")

    # 7. 판매자 과정
    add_title_bullets(prs, "판매자 판매과정 5단계 / 8단계", [
        "저관여 = 바로 쇼핑검색 → 구매",
        "블로그/카페 마케팅하면 헛수고",
        "고관여 = 정보 탐색 → 비교 → 결정",
        "관여도에 따라 판매자 과정이 다르다",
    ], "판매자의 마케팅 전략도 관여도에 맞춰야 합니다.")

    # 8. 상세페이지 전략
    add_title_bullets(prs, "유입경로에 따라 달라지는 상세페이지 전략", [
        "구매과정 6단계를 키워드화",
        "어느 단계를 타겟해서 마케팅하는지 파악",
        "유입 경로별 상세페이지 구성이 달라야 한다",
    ], "같은 상품도 유입 경로에 따라 상세페이지가 달라집니다.")

    # 9. 마케팅 전략
    add_impact(prs, "상품에 따라\n마케팅이 달라진다", "상품별 마케팅 차이를 설명합니다.")
    add_title_bullets(prs, "상품에 따라 달라지는 마케팅 전략", [
        "도마뱀 1,000만원: 길거리 vs 파충류샵",
        "식품: 누구나 먹으니까 OK",
        "성수동 야구용품? → 안 팔림",
        "이걸 온라인화하면 어떻게 달라지는가",
    ], "상품 특성에 맞는 마케팅 채널 선택의 중요성입니다.")

    # 10. Q&A 반박제거
    add_impact_red(prs, "의심과 반박을\nQ&A로 제거하라", "소비자 설득 전략입니다.")
    add_title_bullets(prs, "소비자 가스라이팅", [
        "가격이 비싸면 → 의심이 생긴다",
        "Q&A로 미리 반박을 제거",
        "소비자의 의심을 선제적으로 해소",
        "상세페이지에 FAQ 섹션 필수",
    ], "의심을 제거하면 구매 전환율이 올라갑니다.")

    # 11. 썸네일 제작
    add_section_divider(prs, "실무", "콘텐츠 제작")
    add_impact(prs, "클릭률 높은\n썸네일 만들기", "썸네일 제작 실습입니다.")
    add_title_bullets(prs, "썸네일 제작 실습", [
        "썸네일 공식: 숫자 + 혜택 + 비주얼",
        "나노바나나 / 미리캔버스 활용",
        "실제로 함께 제작",
        "클릭을 부르는 첫인상 만들기",
    ], "실제 툴을 사용하여 썸네일을 만듭니다.")

    # 12. 상세페이지 구성
    add_impact(prs, "상세페이지 구성법", "상세페이지 제작 실습입니다.")
    add_title_bullets(prs, "상세페이지 제작 실습", [
        "저관여: 짧고 핵심 / 고관여: 스토리 + 상세",
        "후킹 → 공감 → 솔루션 → 스펙 → 리뷰 → CTA",
        "수강생 상품으로 실제 제작",
    ], "관여도에 맞는 상세페이지를 함께 만듭니다.")

    # 13. 바이럴 글쓰기
    add_impact(prs, "카페/블로그\n바이럴 글쓰기", "바이럴 글쓰기 실습입니다.")
    add_title_bullets(prs, "바이럴 글쓰기 실습", [
        "키워드 자연스럽게 배치",
        "광고처럼 안 보이게 노출",
        "삭제 안 당하는 팁",
        "실제로 1개 같이 작성",
    ], "실제로 1개를 함께 작성하면서 감을 잡습니다.")

    # 14. 경쟁자 파이
    add_impact_red(prs, "다른 브랜드 키워드를\n내 이득과 연결", "경쟁자의 트래픽을 가져오는 전략입니다.")
    add_title_bullets(prs, "경쟁자 파이 뺏어오는 법", [
        "위고비 가격 검색량 활용",
        "성분 비슷 + 가격 절반 → 설득력",
        "카페 외부형식 실전 적용",
        "경쟁자의 검색 트래픽을 내 것으로",
    ], "경쟁자 키워드를 활용한 마케팅 전략입니다.")

    # 15. 미션
    add_section_divider(prs, "3주차 미션", "")
    add_mission_slide(prs, "3주차 미션 안내", [
        "상세페이지 완성 (D+3)",
        "바이럴 글 2개 작성 (D+5)",
        "경쟁자 파이 뺏기 글 1개 (D+5)",
    ], "3주차 미션을 안내합니다.")

    add_impact(prs, "다음 주에 만나요! 🔥", "3주차를 마무리합니다.")

    prs.save('/Users/kwoneren/.openclaw/workspace/셀러부트캠프_3주차.pptx')
    print(f"Week 3: {len(prs.slides)} slides")

# ============================================================
# WEEK 4
# ============================================================
def gen_week4():
    prs = create_prs()
    add_cover(prs, 4, "스케일업 & 마무리")

    # 1. 시간사용법
    add_section_divider(prs, "OPENING", "시간 사용법")
    add_impact(prs, "시간은 누구에게나\n똑같이 흐를까?", "시간 관리의 중요성을 이야기합니다.")
    add_title_bullets(prs, "남들보다 효율적으로! 시간 사용법", [
        "4주 지난 지금: 누구는 첫매출, 누구는 스토어도 미개설",
        "시간보다 '오늘 뭘 했는지'가 중요",
        "목표지향적 하루를 보내라",
    ], "같은 시간이 주어져도 결과가 다른 이유를 설명합니다.")

    # 2. 친구는 없다
    add_impact_red(prs, "친구는 당분간 없다", "집중을 위한 희생을 이야기합니다.")
    add_title_bullets(prs, "모든 일은 목표지향적으로!", [
        "방해금지 모드 ON, 전화 무시",
        "가족 밥 OK / 친구 술·여행 NO",
        "돈 벌겠다면서 놀러 다니는 건 모순",
        "당분간의 희생이 미래의 자유를 만든다",
    ], "집중의 시기에는 과감한 선택이 필요합니다.")

    # 3. 미션피드백
    add_section_divider(prs, "미션 피드백", "3주차 미션")
    add_title_bullets(prs, "3주차 미션 피드백", [
        "상세페이지 + 바이럴 글 즉석 리뷰",
        "잘한 사례 공유",
        "삭제당한 글 원인 분석",
    ], "3주차 미션 결과를 피드백합니다.")

    # 4. Q&A
    add_section_divider(prs, "Q&A", "셀러들의 궁금증")
    add_impact(prs, "걱정이 앞서\n실행이 안 되는 분들께", "자주 묻는 질문에 답합니다.")
    add_title_bullets(prs, "이 상품 안 팔리면 어떡하죠?", [
        "경쟁사 신고? → 오히려 좋아! 위협됐다는 증거",
        "안 팔리면? → 다음 상품! 재고 없잖아",
        "먼저 걱정해도 나아지는 게 없다",
        "실행이 답이다",
    ], "불안감을 해소하고 실행을 독려합니다.")

    # 5. 수강생 성과
    add_section_divider(prs, "동기부여", "수강생 성과 공유")
    add_impact(prs, "이게 되는구나!", "수강생 성과를 공유합니다.")
    add_title_bullets(prs, "수강생 성과 공유", [
        "4주 동안 매출 나온 수강생 사례 발표",
        "이게 되는구나를 직접 보여주기",
        "아직 안 나온 사람 → 즉석 진단",
    ], "실제 성과를 보여주며 동기를 부여합니다.")

    # 6. 필요 없는 소비자
    add_section_divider(prs, "트레이닝", "")
    add_impact_red(prs, "필요 없는 소비자한테\n파는 법", "니즈 창출 전략입니다.")
    add_title_bullets(prs, "필요로 하지 않는 소비자에게 판매하는 법", [
        "지금 필요 없지만",
        "필요성을 느끼게 해주고 파는 방법",
        "니즈를 만들어내는 것이 진짜 마케팅",
    ], "니즈가 없는 소비자에게도 팔 수 있는 전략입니다.")

    # 7. 제로의 심리
    add_section_divider(prs, "학습", "제로의 심리")
    add_impact(prs, "제로의 심리\n5가지 실전 적용", "심리학 기반 마케팅을 배웁니다.")
    add_title_bullets(prs, "제로의 심리 실전 워크숍", [
        "1. 당위성 — 사야 할 이유 만들기",
        "2. 결핍성 — 한정, 희소성",
        "3. 기발성 — 눈에 띄는 차별화",
        "4. 시급성 — 지금 안 사면 손해",
        "5. 비교성 — 경쟁 상품 대비 우위",
    ], "5가지 심리를 수강생 상품에 직접 적용하는 워크숍입니다.")
    add_title_bullets(prs, "즉석 워크숍", [
        "내 상품에 결핍성을 넣으려면?",
        "시급성을 만들려면?",
        "수강생 상품으로 직접 적용",
    ], "이론이 아닌 실전 적용입니다.")

    # 8. 라면 에피소드
    add_impact(prs, "내가 원한 건 라면\n엄마가 원한 건\n건강한 아들", "소비자 관점의 중요성입니다.")
    add_title_bullets(prs, "그들이 원하는 걸 주고 돈벌기", [
        "소비자가 원하는 게 옳지 않더라도",
        "그걸 줘야 돈이 된다",
        "판매자의 가치관 ≠ 소비자의 욕구",
        "소비자 중심으로 사고하라",
    ], "소비자가 원하는 것을 제공해야 매출이 나옵니다.")

    # 9. 제육 vs 떡볶이
    add_impact_red(prs, "나한테 중요한 건\n중요한 게 아니야", "타겟 중심 사고를 강조합니다.")
    add_title_bullets(prs, "타겟의 취향으로 팔아라", [
        "나는 제육볶음 좋아하지만",
        "여대 앞이면 마라탕/떡볶이",
        "내 취향이 아니라 타겟의 취향",
        "타겟 중심 사고가 매출을 만든다",
    ], "내 취향이 아닌 타겟의 취향으로 판매하라는 메시지입니다.")

    # 10. 검색광고
    add_section_divider(prs, "실무", "광고 & 운영")
    add_impact(prs, "네이버 쇼핑\n검색광고 기초", "광고 실습을 시작합니다.")
    add_title_bullets(prs, "네이버 쇼핑 검색광고 기초", [
        "계정 세팅 → 쇼핑검색 광고",
        "입찰가 → 예산 (하루 5,000원부터)",
        "ROAS 개념 이해",
        "광고 없이도 되지만, 걸면 속도↑",
    ], "검색광고의 기초를 실습합니다.")

    # 11. 주문처리 + CS
    add_title_bullets(prs, "주문처리 + CS 대응", [
        "위탁 주문 → 도매처 발주 프로세스",
        "CS 템플릿: 배송 / 교환 / 환불",
        "리뷰 관리법",
        "톡톡 활용",
    ], "실제 운영에 필요한 주문처리와 CS 대응법입니다.")

    # 12. 스케일업
    add_impact_red(prs, "스케일업 로드맵", "성장 전략을 제시합니다.")
    add_title_bullets(prs, "스케일업 로드맵", [
        "상품 확장: 10 → 50 → 100개",
        "잘 팔리는 패턴으로 유사 확장",
        "쿠팡 추가 진출 (다음 퍼널)",
        "외주화 / 시스템화",
    ], "지속적인 성장을 위한 로드맵을 제시합니다.")

    # 13. Value Ladder
    add_section_divider(prs, "학습", "Value Ladder")
    add_impact(prs, "고객 1명에게\n여러 번 팔아라", "Value Ladder 개념입니다.")
    add_title_bullets(prs, "Value Ladder", [
        "무료 → 저가 → 고가 → 정기구매",
        "고객 1명 확보비용은 같은데",
        "객단가를 올리는 법",
        "여러분 스토어에도 적용 가능",
    ], "고객 생애가치를 높이는 전략입니다.")

    # 14. 마무리
    add_section_divider(prs, "마무리", "4주간의 여정")
    add_impact(prs, "4주간 성과 발표\n+ 다짐", "수강생 성과 발표 시간입니다.")
    add_title_bullets(prs, "수강생 성과 발표 + 마무리", [
        "수강생 성과 발표",
        "첫매출 무한피드백 시작 안내",
        "2기 모집 도우미 혜택 (추천인 할인)",
    ], "4주 과정을 마무리하고 다음 단계를 안내합니다.")

    # 15. 최종 미션
    add_section_divider(prs, "최종 미션", "첫 매출까지 무한 피드백")
    add_mission_slide(prs, "4주차 최종 미션", [
        "첫 매출 인증 (무기한)",
        "상품 5개 이상 등록 (2주 내)",
        "월 목표 + 실행 계획서 (D+3)",
    ], "마지막 미션을 안내합니다.")

    add_impact_red(prs, "첫 매출까지\n무한 피드백!", "과정이 끝나도 지원은 계속됩니다.")
    add_impact(prs, "여러분의 성공을\n응원합니다! 🔥", "마지막 인사입니다.")

    prs.save('/Users/kwoneren/.openclaw/workspace/셀러부트캠프_4주차.pptx')
    print(f"Week 4: {len(prs.slides)} slides")

if __name__ == "__main__":
    gen_week1()
    gen_week2()
    gen_week3()
    gen_week4()
    print("All 4 files generated!")
